#!/usr/bin/env python3
"""Build the updated defense deck: humanize text, indigo rebrand, repurpose slide 13
as Continuous Learning, embed the demo video. Image byte-swap is done separately."""
import copy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

SRC = "/opt/ai-anti-spam-shield-crm/edit-slide-defense/ai-scam-shield.-draft_0.0.3.pptx"
ASSETS = "/opt/ai-anti-spam-shield-crm/edit-slide-defense/assets"
FRAMES = "/opt/ai-anti-spam-shield-crm/defense/defense-docs/video-edit/screenshots"
VIDEO  = "/opt/ai-anti-spam-shield-crm/defense/defense-docs/video-edit/AI_Shield_Demo_Trimmed.mp4"

INDIGO   = RGBColor(0x81,0x8C,0xF8)   # indigo-light: app primaryLight, reads well on dark
INDIGO_D = RGBColor(0x4F,0x46,0xE5)
WHITE    = RGBColor(0xFF,0xFF,0xFF)
CYAN_HEX = {"59DAE4","5CE1E6","D0FAFF","5CE1E6"}

prs = Presentation(SRC)
slides = prs.slides

# ---------- helpers ----------
def first_run_fmt(tf):
    for p in tf.paragraphs:
        for r in p.runs:
            return dict(name=r.font.name, size=r.font.size, bold=r.font.bold,
                        color=(r.font.color.rgb if r.font.color and r.font.color.type is not None else None))
    return dict(name="Lora", size=Pt(18), bold=None, color=WHITE)

def rewrite(shape, lines, bullet="•  ", color=None, header_color=None):
    """Rewrite a text box. `lines` = list of (text, is_header). Preserves first-run font."""
    tf = shape.text_frame
    fmt = first_run_fmt(tf)
    base_color = color or fmt["color"] or WHITE
    tf.clear()
    for i,(text,is_head) in enumerate(lines):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        r = p.add_run()
        r.text = text
        r.font.name = fmt["name"]
        r.font.size = Pt(20) if is_head else (fmt["size"] or Pt(18))
        r.font.bold = True if is_head else fmt["bold"]
        r.font.color.rgb = (header_color or INDIGO) if is_head else base_color

def set_run(shape, newtext):
    """Replace the entire visible text of a single-run-ish box, keep first run fmt."""
    tf = shape.text_frame
    fmt = first_run_fmt(tf)
    tf.clear()
    p = tf.paragraphs[0]; r = p.add_run(); r.text = newtext
    r.font.name=fmt["name"]; r.font.size=fmt["size"]; r.font.bold=fmt["bold"]
    if fmt["color"]: r.font.color.rgb = fmt["color"]

def by_id(slide, sid):
    for sh in slide.shapes:
        if sh.shape_id==sid: return sh
    return None

# ============================================================
# 1) GLOBAL: recolor every cyan text run -> indigo-light
# ============================================================
recolored=0
for s in slides:
    for sh in s.shapes:
        if not sh.has_text_frame: continue
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                try:
                    if r.font.color and r.font.color.type is not None and str(r.font.color.rgb) in CYAN_HEX:
                        r.font.color.rgb = INDIGO; recolored+=1
                except Exception: pass
        if sh.has_table:
            pass
    # tables
    for sh in s.shapes:
        if sh.has_table:
            for row in sh.table.rows:
                for cell in row.cells:
                    for p in cell.text_frame.paragraphs:
                        for r in p.runs:
                            try:
                                if r.font.color and r.font.color.type is not None and str(r.font.color.rgb) in CYAN_HEX:
                                    r.font.color.rgb = INDIGO; recolored+=1
                            except Exception: pass
print("recolored cyan runs ->", recolored)

# ============================================================
# 2) HUMANIZE / FIX TEXT (targeted, format-preserving)
# ============================================================
def find_text(slide, contains):
    for sh in slide.shapes:
        if sh.has_text_frame and contains.lower() in sh.text_frame.text.lower():
            return sh
    return None

# Slide 6 — Limitations heading + body (keep as Limitations slide)
s6 = slides[5]
h=find_text(s6,"LIMITation And Scope")
if h: set_run(h,"Limitations & Scope")
b=find_text(s6,"Platform: Mobile")
if b: rewrite(b, [
    ("Runs as a mobile app — iOS and Android (built with Flutter)",False),
    ("Made for everyday people, not big company setups",False),
    ("Reads English text for now (Khmer is on the roadmap)",False),
    ("Needs internet — the AI runs on our server",False),
])

# Slide 7 — Scope: turn the long wall into simple grouped lines
s7 = slides[6]
h=find_text(s7,"LIMITation And Scope")
if h: set_run(h,"What the App Can Do")
lab=find_text(s7,"Scopes")
if lab: set_run(lab,"Main features")
wall=find_text(s7,"Core scam detection")
if wall: rewrite(wall, [
    ("Spam check — spots junk in SMS and text messages",False),
    ("Voice check — turns a voice message into words, then checks it",False),
    ("Phishing check — catches bad links and fake brand websites",False),
    ("Email auto-scan — watches Gmail, Outlook and Yahoo inboxes",False),
    ("Security dashboard — threats, alerts and network activity in one place",False),
    ("Gets smarter — your feedback retrains the model over time",False),
])

# Slide 8 — fix "Naive Bayes,SVM,"
s8=slides[7]
for sh in s8.shapes:
    if sh.has_table:
        for row in sh.table.rows:
            for cell in row.cells:
                if "Naive Bayes,SVM," in cell.text:
                    for p in cell.text_frame.paragraphs:
                        for r in p.runs:
                            r.text=r.text.replace("Naive Bayes,SVM,","Naive Bayes, SVM")

# Slide 9 — database label
s9=slides[8]
for sh in s9.shapes:
    if sh.has_table:
        for row in sh.table.rows:
            for cell in row.cells:
                if "SQLite / PostgreSQL" in cell.text:
                    for p in cell.text_frame.paragraphs:
                        for r in p.runs:
                            r.text=r.text.replace("SQLite / PostgreSQL with Prisma ORM","PostgreSQL (SQLite for local dev), Prisma ORM")

# Slide 11 — friendlier intro + split note
s11=slides[10]
intro=find_text(s11,"We utilized three specialized")
if intro: set_run(intro,"We trained on three small, clean, balanced datasets from HuggingFace — quality over size.")
# add split note as a new small textbox
tb=s11.shapes.add_textbox(Inches(1.4),Inches(9.3),Inches(17),Inches(0.5))
r=tb.text_frame.paragraphs[0].add_run()
r.text="Split: 80% train / 20% test  ·  stratified  ·  random seed 42 (so results repeat)"
r.font.name="Lora"; r.font.size=Pt(14); r.font.color.rgb=RGBColor(0x94,0xA3,0xB8)

# Slide 12 — phishing content type
s12=slides[11]
for sh in s12.shapes:
    if sh.has_table:
        for row in sh.table.rows:
            for cell in row.cells:
                if cell.text.strip()=="Emails & URLs":
                    for p in cell.text_frame.paragraphs:
                        for r in p.runs:
                            r.text=r.text.replace("Emails & URLs","Website URLs")

# ============================================================
# 3) REPURPOSE SLIDE 13 -> CONTINUOUS LEARNING (+ picture)
# ============================================================
s13=slides[12]
sub=by_id(s13,9)   # "DATASETS"
if sub: set_run(sub,"CONTINUOUS LEARNING")
head=by_id(s13,10) # "Data Split Strategy"
if head: set_run(head,"The app keeps getting smarter")
body=by_id(s13,11) # split text -> short caption
if body: rewrite(body, [
    ("Every time a user confirms or corrects a result, that feedback flows back, "
     "gets checked, and retrains the model — so detection improves on its own.",False),
])
# place continuous-learning diagram
cl = f"{ASSETS}/continuous-learning-indigo.png"
from PIL import Image as PILImage
w,h = PILImage.open(cl).size
pic_w = Inches(13.2); pic_h = Inches(13.2*h/w)
pic_x = (prs.slide_width - pic_w)//2
pic_y = Inches(4.0)
s13.shapes.add_picture(cl, pic_x, pic_y, pic_w, pic_h)

# ============================================================
# 4) SLIDE 28 -> embed demo video + captions
# ============================================================
s28=slides[27]
EMU=914400
sw=prs.slide_width; sh_=prs.slide_height
# title
t=s28.shapes.add_textbox(Inches(10.0),Inches(1.3),Inches(8.6),Inches(1.2))
rp=t.text_frame.paragraphs[0]; rr=rp.add_run(); rr.text="Live Demo"
rr.font.name="Mokoto"; rr.font.size=Pt(40); rr.font.bold=True; rr.font.color.rgb=INDIGO
# caption bullets
cap=s28.shapes.add_textbox(Inches(10.0),Inches(2.8),Inches(8.6),Inches(6.0))
caps=[
 "A real 3-minute walkthrough on a phone:",
 "",
 "Scan a scam SMS  →  flagged 99% spam, CRITICAL",
 "Check a phishing link  →  blocked as very high risk",
 "Auto-scan an email inbox  →  bad mail caught",
 "Report a threat  →  saved to history & dashboard",
]
tf=cap.text_frame; tf.word_wrap=True
for i,line in enumerate(caps):
    p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
    r=p.add_run(); r.text=line
    r.font.name="Lora"; r.font.size=Pt(22 if i==0 else 20)
    r.font.color.rgb = WHITE if i==0 else RGBColor(0xCB,0xD5,0xE1)
# the portrait video
vid_h=Inches(8.8); vid_w=Inches(8.8*1242/2688)  # keep portrait ratio
vid_x=Inches(2.2); vid_y=Inches(1.4)
poster=f"{FRAMES}/06_dashboard_quick_actions.jpg"
s28.shapes.add_movie(VIDEO, vid_x, vid_y, vid_w, vid_h,
                     poster_frame_image=poster, mime_type="video/mp4")

prs.save(SRC)
print("saved deck:", SRC)
print("slide count:", len(prs.slides.__iter__.__self__._sldIdLst))
